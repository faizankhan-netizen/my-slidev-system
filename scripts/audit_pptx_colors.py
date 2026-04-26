import collections 
import collections.abc
from pptx import Presentation

def audit_colors(pptx_path):
    prs = Presentation(pptx_path)
    fonts = []
    
    for slide in prs.slides:
        try:
            # Check for background fill color
            bg = slide.background.fill
            if bg.type == 1: # Solid
                color = bg.fore_color.rgb
                print(f"Slide {prs.slides.index(slide)} BG: #{color}")
            else:
                print(f"Slide {prs.slides.index(slide)} BG: {bg.type}")
        except Exception as e:
            print(f"Slide {prs.slides.index(slide)} BG: DEFAULT (White)")

        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts.append(run.font.name)
            except: pass
                            
    font_counts = collections.Counter(fonts)
    print("\nTOP FONTS DETECTED:")
    for font, count in font_counts.most_common(5):
        print(f"{font}: {count} occurrences")

if __name__ == "__main__":
    audit_colors("AI_Superpower_Enhanced.pptx")
